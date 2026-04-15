"""PPTX assembly — creates the final .pptx file from SVG slides.

Uses python-pptx to create a base presentation, then injects
converted slide XML and media into the PPTX zip package.
"""

from __future__ import annotations

import os
import shutil
import zipfile
import json
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from reportlab.graphics import renderPM
from svglib.svglib import svg2rlg

from backend.config import CANVAS_FORMATS
from backend.generator.svg_finalize.render_ready import prepare_svg_file_for_render

from .converter import convert_svg_to_slide_shapes
from .utils import px_to_emu


def create_pptx(
    svg_files: list[Path],
    output_path: Path,
    canvas_format: str = "ppt169",
    notes: dict[str, str] | None = None,
) -> Path:
    """Create a PPTX file from SVG slides with native DrawingML shapes.

    Args:
        svg_files: Sorted list of SVG file paths.
        output_path: Path for the output .pptx file.
        canvas_format: Canvas format key.
        notes: Dict mapping SVG stem to notes markdown.

    Returns:
        Path to the created PPTX file.
    """
    if notes is None:
        notes = {}

    fmt = CANVAS_FORMATS.get(canvas_format, CANVAS_FORMATS["ppt169"])
    width_emu = Emu(fmt["width"] * 9525)
    height_emu = Emu(fmt["height"] * 9525)

    # Create base PPTX with python-pptx
    prs = Presentation()
    prs.slide_width = width_emu
    prs.slide_height = height_emu

    # Add blank slides (one per SVG)
    blank_layout = prs.slide_layouts[6]  # Blank layout
    for _ in svg_files:
        prs.slides.add_slide(blank_layout)

    # Save to temp file
    tmp_dir = output_path.parent / f".__pptx_build_{uuid.uuid4().hex}"
    tmp_dir.mkdir(parents=True, exist_ok=True)
    try:
        base_pptx = tmp_dir / "base.pptx"
        prs.save(str(base_pptx))

        # Extract PPTX zip
        extract_dir = tmp_dir / "pptx_contents"
        with zipfile.ZipFile(base_pptx, "r") as zf:
            zf.extractall(extract_dir)

        # Ensure media directory exists
        media_dir = extract_dir / "ppt" / "media"
        media_dir.mkdir(parents=True, exist_ok=True)

        # Convert each SVG and inject into the PPTX
        all_media_types: set[str] = set()
        conversion_report: list[dict[str, str | int]] = []
        for i, svg_path in enumerate(svg_files):
            slide_num = i + 1
            prepared_svg_path = prepare_svg_file_for_render(svg_path)
            try:
                try:
                    slide_xml, media_files, rel_entries = convert_svg_to_slide_shapes(
                        prepared_svg_path, slide_num
                    )
                    conversion_report.append({
                        "slide": slide_num,
                        "svg": svg_path.name,
                        "mode": "native",
                    })
                except Exception as exc:
                    slide_xml, media_files, rel_entries = _build_svg_fallback_slide(
                        prepared_svg_path,
                        slide_num,
                        fmt["width"],
                        fmt["height"],
                    )
                    conversion_report.append({
                        "slide": slide_num,
                        "svg": svg_path.name,
                        "mode": "fallback_image",
                        "error_type": type(exc).__name__,
                        "error": str(exc),
                    })
            finally:
                prepared_svg_path.unlink(missing_ok=True)

            # Write slide XML
            slide_file = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            slide_file.write_text(slide_xml, encoding="utf-8")

            # Write media files
            for media_name, media_data in media_files.items():
                media_path = media_dir / media_name
                media_path.write_bytes(media_data)
                ext = media_name.rsplit(".", 1)[-1]
                all_media_types.add(ext)

            # Write relationships
            if rel_entries:
                rels_file = extract_dir / "ppt" / "slides" / "_rels" / f"slide{slide_num}.xml.rels"
                rels_file.parent.mkdir(parents=True, exist_ok=True)
                rels_xml = _build_rels_xml(rel_entries)
                rels_file.write_text(rels_xml, encoding="utf-8")

            # Add speaker notes
            stem = svg_path.stem
            if stem in notes:
                _add_notes(extract_dir, slide_num, notes[stem])

        # Update [Content_Types].xml with media types
        _update_content_types(extract_dir, all_media_types)

        # Repackage as PPTX
        output_path.parent.mkdir(parents=True, exist_ok=True)
        _zip_pptx(extract_dir, output_path)
        report_path = output_path.with_suffix(".conversion_report.json")
        report_path.write_text(json.dumps(conversion_report, ensure_ascii=False, indent=2), encoding="utf-8")

        return output_path
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def _build_rels_xml(rel_entries: list[dict]) -> str:
    """Build slide relationships XML."""
    rels = [
        '<Relationship Id="rId1"'
        ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"'
        ' Target="../slideLayouts/slideLayout1.xml"/>',
    ]
    for entry in rel_entries:
        rels.append(
            f'<Relationship Id="{entry["id"]}" Type="{entry["type"]}"'
            f' Target="{entry["target"]}"/>'
        )

    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        f'{"".join(rels)}'
        '</Relationships>'
    )


def _build_svg_fallback_slide(
    svg_path: Path,
    slide_num: int,
    width_px: int,
    height_px: int,
) -> tuple[str, dict[str, bytes], list[dict]]:
    """Fallback to embedding the whole slide SVG as a single image.

    This keeps export working even when native shape conversion fails.
    """
    png_bytes = _rasterize_svg_to_png(svg_path)
    media_name = f"slide{slide_num}_fallback.png"
    rel_id = "rId2"
    slide_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr>'
        '<a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm>'
        '</p:grpSpPr>'
        '<p:pic>'
        '<p:nvPicPr><p:cNvPr id="2" name="Fallback SVG"/>'
        '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        '<p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{rel_id}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        '<p:spPr>'
        f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="{px_to_emu(width_px)}" cy="{px_to_emu(height_px)}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '</p:spPr>'
        '</p:pic>'
        '</p:spTree></p:cSld>'
        '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>'
        '</p:sld>'
    )
    rel_entries = [{
        "id": rel_id,
        "type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        "target": f"../media/{media_name}",
    }]
    return slide_xml, {media_name: png_bytes}, rel_entries


def _rasterize_svg_to_png(svg_path: Path) -> bytes:
    """Rasterize an SVG to PNG bytes for robust PowerPoint fallback rendering."""
    drawing = svg2rlg(str(svg_path))
    if drawing is None:
        raise ValueError(f"Failed to load SVG drawing for fallback: {svg_path.name}")
    return renderPM.drawToString(drawing, fmt="PNG")


def _add_notes(extract_dir: Path, slide_num: int, notes_text: str) -> None:
    """Add speaker notes to a slide."""
    notes_dir = extract_dir / "ppt" / "notesSlides"
    notes_dir.mkdir(parents=True, exist_ok=True)

    # Simple plaintext notes (strip markdown)
    import re
    plain = re.sub(r"[#*_`\[\]]", "", notes_text)
    plain = re.sub(r"^- ", "• ", plain, flags=re.MULTILINE)

    paragraphs = []
    for line in plain.split("\n"):
        line = line.strip()
        if line:
            escaped = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            paragraphs.append(f'<a:p><a:r><a:t>{escaped}</a:t></a:r></a:p>')
        else:
            paragraphs.append('<a:p><a:endParaRPr/></a:p>')

    notes_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr/>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="Notes"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>'
        '<p:spPr/>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/>{"".join(paragraphs)}</p:txBody>'
        '</p:sp></p:spTree></p:cSld></p:notes>'
    )

    notes_file = notes_dir / f"notesSlide{slide_num}.xml"
    notes_file.write_text(notes_xml, encoding="utf-8")

    # Add relationship from slide to notes
    rels_dir = notes_dir / "_rels"
    rels_dir.mkdir(exist_ok=True)
    notes_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        f'<Relationship Id="rId1"'
        f' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"'
        f' Target="../slides/slide{slide_num}.xml"/>'
        '</Relationships>'
    )
    (rels_dir / f"notesSlide{slide_num}.xml.rels").write_text(notes_rels, encoding="utf-8")


def _update_content_types(extract_dir: Path, media_types: set[str]) -> None:
    """Add media content types to [Content_Types].xml."""
    ct_path = extract_dir / "[Content_Types].xml"
    if not ct_path.exists():
        return

    content = ct_path.read_text(encoding="utf-8")

    ext_map = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "svg": "image/svg+xml",
        "webp": "image/webp",
    }

    insert_pos = content.rfind("</Types>")
    if insert_pos < 0:
        return

    additions = ""
    for ext in media_types:
        mime = ext_map.get(ext)
        if mime and f'Extension="{ext}"' not in content:
            additions += f'<Default Extension="{ext}" ContentType="{mime}"/>'

    if additions:
        content = content[:insert_pos] + additions + content[insert_pos:]
        ct_path.write_text(content, encoding="utf-8")


def _zip_pptx(source_dir: Path, output_path: Path) -> None:
    """Repackage directory as a PPTX zip file."""
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(source_dir):
            for file in files:
                file_path = Path(root) / file
                arc_name = file_path.relative_to(source_dir)
                zf.write(file_path, arc_name)
