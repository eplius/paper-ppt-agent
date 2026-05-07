"""Tests for the PPTX font editor."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Pt

from backend.generator.pptx_font_editor import FontReplaceConfig, replace_fonts_in_pptx


@pytest.fixture
def sample_pptx(workspace_tmp: Path) -> Path:
    """Create a PPTX with text of different sizes for font replacement testing."""
    prs = Presentation()
    prs.slide_width = int(1280 * 9525)
    prs.slide_height = int(720 * 9525)

    blank = prs.slide_layouts[6]

    # Slide 1: heading + body + CJK text
    slide = prs.slides.add_slide(blank)
    txBox = slide.shapes.add_textbox(Pt(100), Pt(100), Pt(800), Pt(400))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Heading (large text)
    p = tf.paragraphs[0]
    p.text = "Research Methodology"
    p.font.size = Pt(28)

    # Body (normal text)
    p = tf.add_paragraph()
    p.text = "This is the body text describing the experimental setup."
    p.font.size = Pt(16)

    # CJK text
    p = tf.add_paragraph()
    p.text = "本研究采用深度学习方法进行实验验证。"
    p.font.size = Pt(18)

    # CJK heading
    p = tf.add_paragraph()
    p.text = "研究方法"
    p.font.size = Pt(28)

    out = workspace_tmp / "test_font_edit.pptx"
    prs.save(str(out))
    return out


def test_replace_fonts_creates_output(sample_pptx: Path, workspace_tmp: Path):
    """Font replacement should produce an output file."""
    config = FontReplaceConfig(
        western_heading="Arial Black",
        western_body="Arial",
        cjk_heading="微软雅黑",
        cjk_body="宋体",
    )
    output = workspace_tmp / "output_fonts.pptx"
    path, result = replace_fonts_in_pptx(sample_pptx, config, output_path=output)

    assert path.exists()
    assert result.fonts_replaced > 0
    assert result.slides_modified > 0


def test_replace_fonts_modifies_xml(sample_pptx: Path, workspace_tmp: Path):
    """Font replacement should change font families in the slide XML."""
    import zipfile

    config = FontReplaceConfig(
        western_heading="Georgia",
        western_body="Times New Roman",
        cjk_heading="SimHei",
        cjk_body="SimSun",
    )
    output = workspace_tmp / "output_fonts2.pptx"
    replace_fonts_in_pptx(sample_pptx, config, output_path=output)

    # Read the modified slide XML
    with zipfile.ZipFile(output, "r") as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")

    # Heading text (28pt) should get Georgia font
    assert "Georgia" in slide_xml
    # Body text (16pt) should get Times New Roman
    assert "Times New Roman" in slide_xml


def test_replace_fonts_partial_config(sample_pptx: Path, workspace_tmp: Path):
    """Partial config (only some fonts set) should work."""
    config = FontReplaceConfig(
        western_body="Georgia",
        # Leave others as None
    )
    output = workspace_tmp / "output_fonts3.pptx"
    path, result = replace_fonts_in_pptx(sample_pptx, config, output_path=output)

    assert path.exists()
    assert result.fonts_replaced > 0  # Should still replace body fonts


def test_replace_fonts_skips_empty(sample_pptx: Path, workspace_tmp: Path):
    """Empty config should produce 0 replacements."""
    config = FontReplaceConfig()  # all None
    output = workspace_tmp / "output_fonts4.pptx"
    path, result = replace_fonts_in_pptx(sample_pptx, config, output_path=output)

    assert path.exists()
    assert result.fonts_replaced == 0


def test_replace_fonts_in_place(sample_pptx: Path, workspace_tmp: Path, monkeypatch):
    """In-place replacement (no output_path) should overwrite the original."""
    import shutil

    # Copy the file so we don't destroy the fixture
    copy = workspace_tmp / "inplace_test.pptx"
    shutil.copy2(sample_pptx, copy)

    config = FontReplaceConfig(western_body="Georgia")
    path, result = replace_fonts_in_pptx(copy, config, output_path=None)

    assert path == copy
    assert result.fonts_replaced > 0
